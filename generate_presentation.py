"""
Generates the Sunflower Solar Panel Controller competition presentation.
Produces:
  - Sunflower_Solar_Controller_Presentation.pptx
  - Sunflower_Solar_Controller_Report.docx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Colour Palette ──────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0xFF, 0xA5, 0x00)   # solar amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE0)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
CARD_BG    = RGBColor(0x15, 0x29, 0x3E)   # slightly lighter navy for cards

# ─── Helpers ─────────────────────────────────────────────────────────────────
def add_slide(prs, layout_index=6):
    layout = prs.slide_layouts[layout_index]   # blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_para(tf, text, font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             space_before=Pt(4), italic=False, bullet_char=""):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if bullet_char:
        run0 = p.add_run()
        run0.text = bullet_char + "  "
        run0.font.size = Pt(font_size)
        run0.font.color.rgb = ACCENT
        run0.font.bold = True
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def add_card(slide, l, t, w, h, title, bullets, title_color=ACCENT,
             bullet_size=11.5, title_size=13):
    add_rect(slide, l, t, w, h, fill_color=CARD_BG, line_color=ACCENT, line_width=Pt(1.2))
    # title
    txb = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.08),
                                   w - Inches(0.24), Inches(0.32))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(title_size)
    run.font.bold  = True
    run.font.color.rgb = title_color
    # bullets
    txb2 = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.40),
                                    w - Inches(0.24), h - Inches(0.50))
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p2 = tf2.paragraphs[0]
            first = False
        else:
            p2 = tf2.add_paragraph()
        p2.space_before = Pt(3)
        r0 = p2.add_run(); r0.text = "\u2022  "; r0.font.color.rgb = ACCENT
        r0.font.size = Pt(bullet_size)
        r = p2.add_run(); r.text = b
        r.font.size = Pt(bullet_size); r.font.color.rgb = WHITE

# ─── Slide definitions ───────────────────────────────────────────────────────
def make_pptx(output_path):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ── SLIDE 1 — Title ──────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)

    # amber top stripe
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)
    # amber bottom stripe
    add_rect(sl, 0, H - Inches(0.07), W, Inches(0.07), ACCENT)

    # central glow circle (decorative)
    add_rect(sl, Inches(9.5), Inches(1.2), Inches(3.5), Inches(3.5),
             fill_color=RGBColor(0x22, 0x3A, 0x52))

    add_text(sl, "SUNFLOWER", Inches(0.6), Inches(1.5), Inches(8), Inches(1.2),
             font_size=52, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    add_text(sl, "Solar Panel Controller", Inches(0.6), Inches(2.7), Inches(9), Inches(0.9),
             font_size=34, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(sl, "An Intelligent, Autonomous Solar Tracking & Monitoring System",
             Inches(0.6), Inches(3.55), Inches(9), Inches(0.6),
             font_size=16, bold=False, color=LIGHT_GREY, align=PP_ALIGN.LEFT, italic=True)

    add_rect(sl, Inches(0.6), Inches(4.25), Inches(4), Inches(0.04), ACCENT)

    add_text(sl, "Competition Presentation  |  April 2026",
             Inches(0.6), Inches(4.45), Inches(7), Inches(0.4),
             font_size=13, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    add_text(sl, "ESP32  \u2022  L298N Motor  \u2022  Servo  \u2022  RTC  \u2022  WiFi TCP  \u2022  OLED",
             Inches(0.6), Inches(5.0), Inches(10), Inches(0.4),
             font_size=12, color=RGBColor(0x88,0xAA,0xCC), align=PP_ALIGN.LEFT, italic=True)

    # ── SLIDE 2 — The Problem ─────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "THE PROBLEM", Inches(0.55), Inches(0.18), Inches(10), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Why fixed solar panels are not enough",
             Inches(0.55), Inches(0.72), Inches(10), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    problems = [
        ("25 – 40%", "Energy Lost",
         "Fixed panels miss peak sun because they cannot track the sun's movement across the sky. A panel pointed east loses hours of afternoon output."),
        ("Manual Only", "No Automation",
         "Repositioning panels by hand is impractical, especially on rooftops, remote farms, or large-scale solar fields with hundreds of panels."),
        ("Zero Visibility", "No Monitoring",
         "Most DIY systems have no way to know current power output, battery health, or whether the system is working — until something fails."),
        ("Weather Damage", "No Protection",
         "Panels left fully exposed during storms, hail, or high wind events are at risk of physical damage. Automated retraction prevents this."),
        ("High Cost", "Barrier to Entry",
         "Commercial solar trackers cost thousands of dollars, making them inaccessible for small farms, rural homes, and developing regions."),
        ("No Remote View", "Operator Blind",
         "Installers and owners have no remote dashboard to see live voltage, current, and power — making fault diagnosis slow and expensive."),
    ]

    cols = 3
    cw = Inches(4.0); ch = Inches(1.75)
    gx = Inches(0.22); gy = Inches(0.20)
    ox = Inches(0.45); oy = Inches(1.25)
    for i, (stat, title, desc) in enumerate(problems):
        col = i % cols; row = i // cols
        x = ox + col * (cw + gx); y = oy + row * (ch + gy)
        add_rect(sl, x, y, cw, ch, CARD_BG, ACCENT, Pt(1))
        add_text(sl, stat, x + Inches(0.15), y + Inches(0.10), cw, Inches(0.45),
                 font_size=22, bold=True, color=ACCENT)
        add_text(sl, title, x + Inches(0.15), y + Inches(0.50), cw, Inches(0.32),
                 font_size=13, bold=True, color=WHITE)
        add_text(sl, desc, x + Inches(0.15), y + Inches(0.82), cw - Inches(0.25), Inches(0.88),
                 font_size=10.5, color=LIGHT_GREY, wrap=True)

    add_rect(sl, Inches(0.55), H - Inches(0.55), Inches(12.2), Inches(0.03), ACCENT)
    add_text(sl, "Our system solves all six problems with a single $30 device.",
             Inches(0.55), H - Inches(0.50), Inches(12), Inches(0.38),
             font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ── SLIDE 3 — What it Does ────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "WHAT THE SYSTEM DOES", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Inspired by the sunflower — follows the sun from dawn to dusk",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    features = [
        ("OPEN",      "Opens panel at sunrise",     "RTC-scheduled DC motor fires at programmed morning time, physically deploying the solar panel to face the sky."),
        ("TRACK",     "Tracks sun every 15 minutes","Continuous-rotation servo rotates the panel angle toward the sun at 15-minute intervals throughout the day."),
        ("CLOSE",     "Closes panel at sunset",     "Motor reverses at programmed evening time, retracting the panel to protect it from wind, rain, and vandalism overnight."),
        ("MONITOR",   "Monitors output live",       "ACS712 current sensor + voltage dividers measure panel voltage (V), current (A), and calculated power (W) every 100 ms."),
        ("DISPLAY",   "Wireless OLED dashboard",    "A separate ESP32 shows live sensor values on an OLED screen over a private WiFi network — no internet required."),
        ("OVERRIDE",  "Manual control anytime",     "4 physical buttons let the operator manually open/close the panel or nudge the servo angle at any time."),
        ("RESET",     "Emergency stop",             "A single RESET command instantly halts motor and servo, clears all state flags — safe for field use."),
        ("CALIBRATE", "Auto sensor calibration",    "On every boot, 500 ADC samples are averaged to auto-calibrate the zero-current offset — no manual tuning needed."),
    ]

    cw = Inches(3.1); ch = Inches(1.55)
    gx = Inches(0.16); gy = Inches(0.14)
    ox = Inches(0.45); oy = Inches(1.22)
    for i, (badge, title, desc) in enumerate(features):
        col = i % 4; row = i // 4
        x = ox + col * (cw + gx); y = oy + row * (ch + gy)
        add_rect(sl, x, y, cw, ch, CARD_BG, ACCENT, Pt(1))
        add_text(sl, badge, x + Inches(0.12), y + Inches(0.08), cw, Inches(0.3),
                 font_size=10, bold=True, color=ACCENT)
        add_text(sl, title, x + Inches(0.12), y + Inches(0.38), cw - Inches(0.2), Inches(0.32),
                 font_size=12.5, bold=True, color=WHITE)
        add_text(sl, desc, x + Inches(0.12), y + Inches(0.68), cw - Inches(0.2), Inches(0.82),
                 font_size=10, color=LIGHT_GREY, wrap=True)

    # ── SLIDE 4 — Architecture ────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "SYSTEM ARCHITECTURE", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Two ESP32 units — fully wireless, no internet required",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    # Controller box
    cx = Inches(0.55); cy = Inches(1.28); cw2 = Inches(4.6); ch2 = Inches(5.5)
    add_rect(sl, cx, cy, cw2, ch2, CARD_BG, ACCENT, Pt(1.5))
    add_text(sl, "Solar Controller  (ESP32 — Client)", cx + Inches(0.15), cy + Inches(0.12),
             cw2, Inches(0.38), font_size=13, bold=True, color=ACCENT)
    ctrl_items = [
        "L298N DC Motor Driver",
        "Continuous Rotation Servo",
        "Panel Voltage Sensor (GPIO 34)",
        "Panel Current Sensor ACS712 (GPIO 35)",
        "Battery Voltage Sensor (GPIO 32)",
        "RTC DS3231 — I\u00B2C timekeeping",
        "WiFi Station mode — TCP client",
        "Non-blocking 10ms loop",
        "PWM soft-ramp motor control",
        "Auto current sensor calibration",
    ]
    txb = sl.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.58),
                                cw2 - Inches(0.3), ch2 - Inches(0.7))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for item in ctrl_items:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.space_before = Pt(4)
        r0 = p.add_run(); r0.text = "\u25B6  "; r0.font.color.rgb = ACCENT; r0.font.size = Pt(10)
        r = p.add_run(); r.text = item; r.font.size = Pt(11); r.font.color.rgb = WHITE

    # Arrow (WiFi)
    mx = Inches(5.35); my = Inches(3.8)
    add_rect(sl, mx, my, Inches(2.5), Inches(0.04), ACCENT)
    add_text(sl, "\u25BA", mx + Inches(2.3), my - Inches(0.18), Inches(0.3), Inches(0.4),
             font_size=14, color=ACCENT)
    add_text(sl, "\u25C4", mx - Inches(0.25), my - Inches(0.18), Inches(0.3), Inches(0.4),
             font_size=14, color=ACCENT)
    add_text(sl, "WiFi TCP", mx + Inches(0.55), my - Inches(0.38), Inches(1.4), Inches(0.35),
             font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, "Sensor CSV (100ms)", mx, my + Inches(0.1), Inches(2.5), Inches(0.3),
             font_size=9.5, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "Button Commands", mx, my + Inches(0.38), Inches(2.5), Inches(0.3),
             font_size=9.5, color=RGBColor(0xFF,0xCC,0x44), align=PP_ALIGN.CENTER)

    # Remote box
    rx = Inches(8.0); ry = Inches(1.28); rw = Inches(4.75); rh2 = Inches(5.5)
    add_rect(sl, rx, ry, rw, rh2, CARD_BG, ACCENT, Pt(1.5))
    add_text(sl, "Remote Control Unit  (ESP32 — AP)", rx + Inches(0.15), ry + Inches(0.12),
             rw, Inches(0.38), font_size=13, bold=True, color=ACCENT)
    remote_items = [
        "WiFi Access Point (192.168.4.1:5000)",
        "OLED 128\u00D764 live dashboard",
        "Displays: Panel V, Current A, Power W",
        "Displays: Battery V, Time (HH:MM:SS)",
        "Button 1 \u2014 SERVO_FWD (hold to run)",
        "Button 2 \u2014 SERVO_REV (hold to run)",
        "Button 3 \u2014 PANEL_OPEN (hold to run)",
        "Button 4 \u2014 PANEL_CLOSE (hold to run)",
        "Auto-reconnects dropped clients",
        "Display refresh every 500ms",
    ]
    txb2 = sl.shapes.add_textbox(rx + Inches(0.15), ry + Inches(0.58),
                                  rw - Inches(0.3), rh2 - Inches(0.7))
    tf2 = txb2.text_frame; tf2.word_wrap = True
    first = True
    for item in remote_items:
        if first: p = tf2.paragraphs[0]; first = False
        else: p = tf2.add_paragraph()
        p.space_before = Pt(4)
        r0 = p.add_run(); r0.text = "\u25B6  "; r0.font.color.rgb = ACCENT; r0.font.size = Pt(10)
        r = p.add_run(); r.text = item; r.font.size = Pt(11); r.font.color.rgb = WHITE

    # ── SLIDE 5 — Key Technical Features ─────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "KEY TECHNICAL FEATURES", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Engineering decisions that make the system reliable in the field",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    tech = [
        ("Non-Blocking Architecture",
         ["Main loop runs every 10 ms — never uses delay()",
          "All operations use millis() timers (motor, servo, WiFi, sensors)",
          "Motor, servo and WiFi reconnect all run concurrently without pausing each other",
          "WiFi connect() is only called when motor is IDLE to prevent blocking the PWM ramp"]),
        ("Motor State Machine",
         ["States: MOTOR_IDLE \u2192 OPENING / CLOSING / MANUAL \u2192 IDLE",
          "New commands only accepted when state is MOTOR_IDLE — prevents double-trigger",
          "opened/closed flags set AFTER motor physically finishes — reflects real world",
          "Emergency RESET command instantly halts motor, servo, and clears all flags"]),
        ("PWM Soft-Ramp Control",
         ["Custom MotorL298N library ramps 0\u219290% over ~5 seconds",
          "30% kickstart PWM prevents stall against static friction on startup",
          "hardBrake() instantly cuts all PWM — used during RESET",
          "Smooth ramp reduces mechanical shock, extends belt and gear life significantly"]),
        ("Automatic Sensor Calibration",
         ["ACS712 zero-current offset auto-calibrated on every boot",
          "Averages 500 ADC samples over 1 second before accepting load",
          "Eliminates offset drift from temperature and supply voltage variation",
          "No manual potentiometer tuning required — works out of the box"]),
        ("Servo Sun Tracking",
         ["Auto-tracking fires every 15 minutes while panel is open",
          "Duration calculated from RPM and revolutions (not guessed)",
          "Manual override has highest priority — auto-track pauses instantly",
          "Tracking cancelled automatically when panel closes at end of day"]),
        ("Data Transmission",
         ["CSV sent over TCP every 100 ms: time, V, A, W, battery V",
          "RTC cached every 1 second — avoids I\u00B2C overhead every 10 ms loop",
          "WiFi auto-reconnects every 5 seconds if connection drops",
          "Button press edge-detected with 200 ms debounce; release is immediate"]),
    ]

    cw3 = Inches(4.0); ch3 = Inches(2.5)
    gx3 = Inches(0.22); gy3 = Inches(0.18)
    ox3 = Inches(0.45); oy3 = Inches(1.22)
    for i, (title, bullets) in enumerate(tech):
        col = i % 3; row = i // 3
        x = ox3 + col * (cw3 + gx3); y = oy3 + row * (ch3 + gy3)
        add_card(sl, x, y, cw3, ch3, title, bullets, bullet_size=10.5)

    # ── SLIDE 6 — Communication Protocol ─────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "COMMUNICATION PROTOCOL", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Private WiFi — TCP — no cloud, no internet, no latency",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    # Controller -> Remote
    add_rect(sl, Inches(0.55), Inches(1.2), Inches(12.2), Inches(2.6), CARD_BG, ACCENT, Pt(1))
    add_text(sl, "Controller  \u2192  Remote  (every 100 ms)",
             Inches(0.75), Inches(1.28), Inches(10), Inches(0.38),
             font_size=14, bold=True, color=GREEN)
    add_text(sl, "Format:   HH:MM:SS , voltage , current , power , batteryVoltage",
             Inches(0.75), Inches(1.65), Inches(11), Inches(0.35),
             font_size=13, bold=True, color=WHITE)
    add_text(sl, "Example:  08:42:15 , 18.34 , 2.10 , 38.51 , 12.67",
             Inches(0.75), Inches(2.0), Inches(11), Inches(0.35),
             font_size=12.5, color=ACCENT, italic=True)

    fields = [
        ("HH:MM:SS", "Time", "DS3231 RTC"),
        ("voltage", "V", "Panel voltage divider GPIO 34"),
        ("current", "A", "ACS712 current sensor GPIO 35"),
        ("power", "W", "Calculated: voltage \u00D7 current"),
        ("batteryVoltage", "V", "Battery voltage divider GPIO 32"),
    ]
    col_x = [Inches(0.75), Inches(3.5), Inches(5.5), Inches(8.5)]
    add_text(sl, "Field", col_x[0], Inches(2.42), Inches(2.6), Inches(0.28), font_size=11, bold=True, color=ACCENT)
    add_text(sl, "Unit",  col_x[1], Inches(2.42), Inches(1.8), Inches(0.28), font_size=11, bold=True, color=ACCENT)
    add_text(sl, "Source", col_x[2], Inches(2.42), Inches(3.5), Inches(0.28), font_size=11, bold=True, color=ACCENT)
    for j, (field, unit, source) in enumerate(fields):
        yy = Inches(2.72) + j * Inches(0.28)
        c = WHITE if j % 2 == 0 else LIGHT_GREY
        add_text(sl, field,  col_x[0], yy, Inches(2.6), Inches(0.28), font_size=11, color=c)
        add_text(sl, unit,   col_x[1], yy, Inches(1.8), Inches(0.28), font_size=11, color=c)
        add_text(sl, source, col_x[2], yy, Inches(4.5), Inches(0.28), font_size=11, color=c)

    # Remote -> Controller
    add_rect(sl, Inches(0.55), Inches(4.0), Inches(12.2), Inches(2.85), CARD_BG, ACCENT, Pt(1))
    add_text(sl, "Remote  \u2192  Controller  (on button press/release)",
             Inches(0.75), Inches(4.08), Inches(10), Inches(0.38),
             font_size=14, bold=True, color=RGBColor(0xFF,0xCC,0x44))
    cmds = [
        ("SERVO_FWD",   "Hold button 1", "Servo rotates forward — manual sun tracking"),
        ("SERVO_REV",   "Hold button 2", "Servo rotates reverse — undo tracking angle"),
        ("SERVO_STOP",  "Release button 1 or 2", "Servo stops immediately on release"),
        ("PANEL_OPEN",  "Hold button 3", "DC motor opens panel (only if MOTOR_IDLE)"),
        ("PANEL_CLOSE", "Hold button 4", "DC motor closes panel (only if MOTOR_IDLE)"),
        ("MOTOR_STOP",  "Release button 3 or 4", "Motor stops immediately on release"),
        ("RESET",       "Emergency", "Halts everything, clears all state flags"),
    ]
    cols_c = [Inches(0.75), Inches(3.8), Inches(7.2)]
    add_text(sl, "Command",  cols_c[0], Inches(4.5), Inches(2.8), Inches(0.28), font_size=11, bold=True, color=ACCENT)
    add_text(sl, "Trigger",  cols_c[1], Inches(4.5), Inches(3.1), Inches(0.28), font_size=11, bold=True, color=ACCENT)
    add_text(sl, "Effect",   cols_c[2], Inches(4.5), Inches(5.5), Inches(0.28), font_size=11, bold=True, color=ACCENT)
    for j, (cmd, trigger, effect) in enumerate(cmds):
        yy = Inches(4.8) + j * Inches(0.285)
        c = WHITE if j % 2 == 0 else LIGHT_GREY
        add_text(sl, cmd,     cols_c[0], yy, Inches(2.8), Inches(0.28), font_size=10.5, color=GREEN)
        add_text(sl, trigger, cols_c[1], yy, Inches(3.1), Inches(0.28), font_size=10.5, color=c)
        add_text(sl, effect,  cols_c[2], yy, Inches(5.5), Inches(0.28), font_size=10.5, color=c)

    # ── SLIDE 7 — Large Scale Implementation ─────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "LARGE SCALE IMPLEMENTATION", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "How this prototype scales to a commercial solar farm",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    scale_cards = [
        ("Hardware Scaling",
         ["Replace ESP32 with Raspberry Pi / industrial PLC as central hub",
          "L298N \u2192 3-phase motor controllers for larger actuators",
          "Single servo \u2192 stepper motors with encoders for precise angular control",
          "Voltage dividers \u2192 industrial CT clamps rated for high currents (100A+)",
          "OLED \u2192 web-based Grafana/React dashboard for fleet-wide metrics",
          "Add dual-axis trackers: separate motor for elevation angle control"]),
        ("Network & Software Scaling",
         ["Replace direct TCP with MQTT broker (Mosquitto) — each panel is a topic",
          "OTA (Over-the-Air) firmware updates across entire field without physical access",
          "Time-series database (InfluxDB / TimescaleDB) for historical analysis",
          "REST API so any device (phone, browser, SCADA) can query panel status",
          "GPS-based sun angle calculation from latitude/longitude using NOAA formulas",
          "Cloud backup with daily energy yield reports sent via email/SMS"]),
        ("Physical Deployment",
         ["Single-axis tracker (east-west) driven by DC motor — most common commercial design",
          "Dual-axis tracker adds elevation control for maximum annual yield",
          "Wind speed sensor auto-retracts all panels when gusts exceed safe threshold",
          "IP65 weatherproof enclosures for all electronics; surge protection on sensors",
          "Modular wiring harness — each panel unit is identical, plug-and-play",
          "Remote diagnostic port: plug in laptop to pull logs without disassembly"]),
        ("Operations at Scale",
         ["Central dashboard shows all panels on a map — green/amber/red status",
          "Predictive maintenance: compare daily yield to expected yield, flag underperformers",
          "Automated fault detection: if panel V drops unexpectedly, send SMS alert",
          "Staggered motor schedules to avoid current surge when 100 panels open at once",
          "Role-based access: operators can view, admins can reconfigure schedules",
          "Audit log of all commands and state changes for insurance and warranty claims"]),
    ]

    cw4 = Inches(6.0); ch4 = Inches(2.75)
    gx4 = Inches(0.22); gy4 = Inches(0.18)
    ox4 = Inches(0.45); oy4 = Inches(1.22)
    for i, (title, bullets) in enumerate(scale_cards):
        col = i % 2; row = i // 2
        x = ox4 + col * (cw4 + gx4); y = oy4 + row * (ch4 + gy4)
        add_card(sl, x, y, cw4, ch4, title, bullets, bullet_size=11)

    # ── SLIDE 8 — Improvements ────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "SYSTEM IMPROVEMENTS", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Upgrades to transform the prototype into a production-grade system",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    imp_cards = [
        ("Precision Sun Tracking",
         ["Add LDR (light-dependent resistor) pair — servo stops when both read equal light",
          "Closed-loop: encoder on servo shaft confirms actual angular position",
          "MPPT algorithm — track angle that maximises real-time power, not just sun direction",
          "Add elevation axis motor for full 2-axis tracking (+15% extra yield over single axis)",
          "Store tracking history to learn seasonal patterns and anticipate sun position"]),
        ("Reliability & Safety",
         ["Motor stall detection — dedicated current sensor on motor supply line",
          "Panel position encoder — know exact open angle, prevent over-driving",
          "ESP32 hardware watchdog — auto-reboot if firmware hangs",
          "Redundant NTP time sync when WiFi is available; RTC as offline fallback",
          "Surge protection and reverse-polarity protection on all power inputs"]),
        ("Energy Intelligence",
         ["State-of-charge calculation from battery voltage discharge curve",
          "Watt-hour accumulation — log daily energy yield (kWh) to flash memory",
          "Low battery alert — automatically retract panel and go to sleep mode",
          "Compare actual vs expected yield — alert if efficiency drops below threshold",
          "Temperature sensor on panel — derate power estimate on hot days (3% per \u00B0C)"]),
        ("User Interface",
         ["Replace OLED with colour TFT touchscreen showing voltage/current/power graphs",
          "Mobile app via Bluetooth for configuration (set open/close times, tracking interval)",
          "Push notifications when power drops below expected or a fault is detected",
          "Voice alerts via small buzzer — different tones for open, close, fault events",
          "Web interface with historical charts, exportable CSV, and remote schedule config"]),
        ("Connectivity",
         ["MQTT over WiFi to a local home server or cloud broker (AWS IoT / HiveMQ)",
          "LoRa radio fallback for installations beyond WiFi range (up to 10km line-of-sight)",
          "4G LTE module option for completely off-grid sites with no WiFi infrastructure",
          "Mesh networking — controllers relay data through each other to the gateway",
          "Encrypted TLS communication to prevent command injection over the air"]),
        ("Manufacturing & Deployment",
         ["Custom PCB to replace breadboard wiring — reduces failure points significantly",
          "Conformal coating on PCB for humidity and corrosion resistance",
          "3D-printed weatherproof enclosure with cable glands and DIN rail mount",
          "Factory calibration routine — automated test jig checks each unit before shipping",
          "QR code on each unit links to its data dashboard and service history"]),
    ]

    cw5 = Inches(4.0); ch5 = Inches(2.65)
    gx5 = Inches(0.22); gy5 = Inches(0.16)
    ox5 = Inches(0.45); oy5 = Inches(1.22)
    for i, (title, bullets) in enumerate(imp_cards):
        col = i % 3; row = i // 3
        x = ox5 + col * (cw5 + gx5); y = oy5 + row * (ch5 + gy5)
        add_card(sl, x, y, cw5, ch5, title, bullets, bullet_size=10.5)

    # ── SLIDE 9 — Future Benefits ─────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "FUTURE BENEFITS & APPLICATIONS", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Why this technology matters beyond the competition",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    future_cards = [
        ("Energy Independence",
         ["Single-axis tracker increases annual yield by 25\u201340% over fixed panels",
          "Dual-axis tracking adds a further 10\u201315% on top of single-axis",
          "Fewer panels needed for the same power target \u2014 lower land and material cost",
          "Longer battery cycle life: tracker extracts full energy during peak sun hours",
          "Enables off-grid homes, farms, and villages to become energy self-sufficient"]),
        ("Agricultural (Agrivoltaics)",
         ["Tilted panels let morning/evening light through while blocking harsh midday sun",
          "Crops grown underneath benefit from shade \u2014 up to 60% less water consumption",
          "Dual income: farmers sell electricity AND grow crops on the same land",
          "This system\u2019s tilt angle control makes it a ready foundation for agrivoltaics",
          "Already deployed at scale in Japan, Germany, and South Korea"]),
        ("Smart Grid Integration",
         ["Real-time power data at 100ms resolution feeds directly into grid balancing",
          "When panels report excess generation, smart switches divert energy to EV charging",
          "Demand-response: utilities can remotely schedule open/close times to shift load",
          "Micro-grid: cluster of these units can power a village with automatic load sharing",
          "Predictive generation forecast helps grid operators plan backup generation"]),
        ("Developing World & Disaster Relief",
         ["No internet required \u2014 the two ESP32s form their own private WiFi network",
          "Total BOM cost under \u20B92500 (\u223C$30) \u2014 accessible to smallholder farmers globally",
          "Can be assembled with basic tools and locally available components",
          "Solar tracker + battery means uninterrupted power even in remote locations",
          "Deployed to disaster zones as portable power stations for medical equipment"]),
        ("Environmental Impact",
         ["Every 1% yield improvement = proportionally fewer panels for a given power target",
          "Fewer panels = less silicon, aluminium, and rare earth mining required",
          "1 MW solar farm with tracking: +350 kWh/day = saves ~50,000 kg CO\u2082/year",
          "Retractable panels survive hailstorms \u2014 longer panel life, less e-waste",
          "Closed-loop MPPT maximises green energy extraction without any fuel burning"]),
        ("Commercial & Research Opportunities",
         ["Platform for AI-based sun prediction using historical tracking data",
          "University research tool for studying panel degradation and soiling effects",
          "Integration with weather APIs to pre-position panels before storms arrive",
          "Data marketplace: anonymised energy yield data valuable to grid planners",
          "White-label product for solar installation companies to brand and resell"]),
    ]

    cw6 = Inches(4.0); ch6 = Inches(2.65)
    gx6 = Inches(0.22); gy6 = Inches(0.16)
    ox6 = Inches(0.45); oy6 = Inches(1.22)
    for i, (title, bullets) in enumerate(future_cards):
        col = i % 3; row = i // 3
        x = ox6 + col * (cw6 + gx6); y = oy6 + row * (ch6 + gy6)
        add_card(sl, x, y, cw6, ch6, title, bullets, bullet_size=10.5)

    # ── SLIDE 10 — Impact Numbers ─────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "REAL-WORLD IMPACT", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Quantified energy gains from single-axis solar tracking",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    stats = [
        ("+35%", "Average annual energy gain\nfrom single-axis tracking\nvs fixed panel"),
        ("+50%", "Peak-hour energy gain\non clear summer days"),
        ("100ms", "Sensor data refresh rate\n— real-time monitoring\nfor every panel"),
        ("\u20B92500", "Total Bill of Materials\ncost for the complete\ntwo-unit system"),
        ("15 min", "Sun tracking interval\n— optimal balance of\nyield vs servo wear"),
        ("50,000 kg", "CO\u2082 saved per year\nfor a 1 MW farm\nusing this tracker"),
    ]
    sw = Inches(4.0); sh = Inches(2.3)
    sgx = Inches(0.22); sgy = Inches(0.22)
    sox = Inches(0.45); soy = Inches(1.3)
    for i, (num, label) in enumerate(stats):
        col = i % 3; row = i // 3
        x = sox + col * (sw + sgx); y = soy + row * (sh + sgy)
        add_rect(sl, x, y, sw, sh, CARD_BG, ACCENT, Pt(1.5))
        add_text(sl, num, x + Inches(0.15), y + Inches(0.18), sw - Inches(0.25), Inches(0.7),
                 font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(sl, label, x + Inches(0.1), y + Inches(0.9), sw - Inches(0.2), Inches(1.3),
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

    add_rect(sl, Inches(0.55), H - Inches(0.6), Inches(12.2), Inches(0.03), ACCENT)
    add_text(sl,
             "Based on 35% improvement from single-axis tracking, 300 sunny days/year, 0.4 kg CO\u2082/kWh grid average (IEA 2024)",
             Inches(0.55), H - Inches(0.52), Inches(12.2), Inches(0.35),
             font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    # ── SLIDE 11 — Competitive Advantages ────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "COMPETITIVE ADVANTAGES", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "What makes this system stand out",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    adv = [
        ("01", "No Internet Required",
         "The two ESP32s form a private WiFi network. Works in fields, deserts, rooftops — anywhere with no infrastructure. Unlike cloud-dependent systems, it cannot be killed by a server outage."),
        ("02", "Custom PWM Motor Library",
         "Soft-ramp 0\u219290% over 5 seconds with kickstart — prevents gear shock, motor burnout, and belt snap. No commercial tracker at this price point offers configurable ramp control."),
        ("03", "Self-Calibrating Sensors",
         "ACS712 zero offset auto-calibrated on every boot with 500 samples. Eliminates manual potentiometer tuning and compensates for temperature drift automatically over the lifetime of the unit."),
        ("04", "State-Machine Safety",
         "Motor only accepts a new command when IDLE — impossible to trigger opening while already opening. Flags are set after physical completion, not on command issue, reflecting real-world state."),
        ("05", "100ms Live Data Stream",
         "10 readings per second over TCP gives a smooth, real-time view of voltage, current, and power. Enough resolution to detect transient shadows, partial shading, and wiring faults."),
        ("06", "Extensible Command Protocol",
         "Adding a new button or remote command requires one line in the buttons[] array — no other code changes. Engineers can add features in minutes without understanding the full codebase."),
        ("07", "Emergency Reset",
         "A single RESET command stops motor, stops servo, clears all flags and timers instantly. Designed for field safety — operator presses one button and the entire system is in a known safe state."),
        ("08", "Under \u20B92500 Total Cost",
         "Complete two-unit system costs under \u20B92500 using off-the-shelf modules. Commercial solar trackers start at \u20B940,000+. This makes professional-grade tracking accessible to small farmers globally."),
    ]

    aw = Inches(5.9); ah = Inches(1.55)
    agx = Inches(0.22); agy = Inches(0.15)
    aox = Inches(0.45); aoy = Inches(1.22)
    for i, (num, title, desc) in enumerate(adv):
        col = i % 2; row = i // 2
        x = aox + col * (aw + agx); y = aoy + row * (ah + agy)
        add_rect(sl, x, y, aw, ah, CARD_BG, ACCENT, Pt(1))
        add_text(sl, num, x + Inches(0.12), y + Inches(0.12), Inches(0.5), Inches(0.4),
                 font_size=18, bold=True, color=ACCENT)
        add_text(sl, title, x + Inches(0.65), y + Inches(0.12), aw - Inches(0.75), Inches(0.38),
                 font_size=13, bold=True, color=WHITE)
        add_text(sl, desc, x + Inches(0.12), y + Inches(0.55), aw - Inches(0.22), Inches(0.95),
                 font_size=10.5, color=LIGHT_GREY, wrap=True)

    # ── SLIDE 12 — Bill of Materials ──────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)

    add_text(sl, "BILL OF MATERIALS", Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
             font_size=28, bold=True, color=ACCENT)
    add_text(sl, "Complete hardware cost — under \u20B92500 for the full system",
             Inches(0.55), Inches(0.72), Inches(11), Inches(0.38),
             font_size=15, color=LIGHT_GREY, italic=True)
    add_rect(sl, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.03), ACCENT)

    bom = [
        ("ESP32 Dev Board",           "2",  "\u20B9400",  "\u20B9800",  "Main controller + Remote control unit"),
        ("L298N Motor Driver Module", "1",  "\u20B9150",  "\u20B9150",  "DC motor direction and PWM speed control"),
        ("DC Gear Motor (12V)",       "1",  "\u20B9350",  "\u20B9350",  "Opens and closes the solar panel"),
        ("Continuous Rotation Servo", "1",  "\u20B9280",  "\u20B9280",  "Sun tracking axis rotation"),
        ("ACS712-30A Current Sensor", "1",  "\u20B9120",  "\u20B9120",  "Panel output current measurement"),
        ("Voltage Divider Module",    "2",  "\u20B945",   "\u20B990",   "Panel voltage + battery voltage sensing"),
        ("DS3231 RTC Module",         "1",  "\u20B9180",  "\u20B9180",  "Real-time clock for open/close scheduling"),
        ("SSD1306 OLED 128\u00D764",  "1",  "\u20B9220",  "\u20B9220",  "Live sensor dashboard on remote unit"),
        ("Push Buttons (pack)",       "1",  "\u20B950",   "\u20B950",   "4 control buttons on remote unit"),
        ("Miscellaneous (wires etc)", "1",  "\u20B9200",  "\u20B9200",  "Jumper wires, breadboard, capacitors, resistors"),
    ]

    # Header
    cols_b = [Inches(0.65), Inches(4.85), Inches(6.2), Inches(7.5), Inches(8.8)]
    add_rect(sl, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.38), ACCENT)
    hdrs = ["Component", "Qty", "Unit Price", "Total", "Purpose"]
    for j, (hdr, cx2) in enumerate(zip(hdrs, cols_b)):
        add_text(sl, hdr, cx2, Inches(1.26), Inches(2), Inches(0.32),
                 font_size=12, bold=True, color=DARK_BG)

    for i, (comp, qty, unit, total, purpose) in enumerate(bom):
        y2 = Inches(1.62) + i * Inches(0.44)
        bg2 = CARD_BG if i % 2 == 0 else RGBColor(0x18,0x30,0x48)
        add_rect(sl, Inches(0.55), y2, Inches(12.2), Inches(0.42), bg2)
        vals = [comp, qty, unit, total, purpose]
        for j, (val, cx2) in enumerate(zip(vals, cols_b)):
            add_text(sl, val, cx2, y2 + Inches(0.06), Inches(2.1), Inches(0.34),
                     font_size=11, color=WHITE if j != 3 else GREEN)

    # Total row
    ty = Inches(1.62) + len(bom) * Inches(0.44)
    add_rect(sl, Inches(0.55), ty, Inches(12.2), Inches(0.44), ACCENT)
    add_text(sl, "TOTAL", cols_b[0], ty + Inches(0.06), Inches(3), Inches(0.34),
             font_size=13, bold=True, color=DARK_BG)
    add_text(sl, "\u20B92,440", cols_b[3], ty + Inches(0.06), Inches(1.5), Inches(0.34),
             font_size=13, bold=True, color=DARK_BG)

    # ── SLIDE 13 — Conclusion ─────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, W, Inches(0.07), ACCENT)
    add_rect(sl, 0, H - Inches(0.07), W, Inches(0.07), ACCENT)

    add_text(sl, "CONCLUSION", Inches(0.6), Inches(0.5), Inches(11), Inches(0.6),
             font_size=34, bold=True, color=ACCENT)
    add_rect(sl, Inches(0.6), Inches(1.18), Inches(5), Inches(0.04), ACCENT)

    points = [
        ("Low-cost hardware, professional results",
         "Under \u20B92500 of components delivers intelligent, autonomous solar tracking — capabilities previously only available in \u20B940,000+ commercial systems."),
        ("Non-blocking, state-machine firmware",
         "Motor, servo, WiFi, and sensors all run concurrently on a single microcontroller. The software architecture is sound enough to scale to production."),
        ("Wireless and self-contained",
         "No internet, no cloud, no single point of failure. The two-unit system forms its own private network — deployable anywhere on Earth."),
        ("Clear upgrade path",
         "LDR closed-loop tracking, MPPT, MQTT, OTA updates, custom PCB — every improvement is well-defined and builds directly on the existing architecture."),
        ("Measurable real-world impact",
         "+35% annual energy yield, 50,000 kg CO\u2082 saved per MW, and a \u20B92500 price point that makes solar tracking accessible to smallholder farmers worldwide."),
    ]
    for i, (title, body) in enumerate(points):
        y3 = Inches(1.35) + i * Inches(1.02)
        add_rect(sl, Inches(0.6), y3, Inches(12.1), Inches(0.96), CARD_BG, ACCENT, Pt(0.8))
        add_text(sl, f"{i+1}.  {title}", Inches(0.78), y3 + Inches(0.06),
                 Inches(11.7), Inches(0.34), font_size=13, bold=True, color=ACCENT)
        add_text(sl, body, Inches(0.78), y3 + Inches(0.40),
                 Inches(11.7), Inches(0.52), font_size=11.5, color=WHITE, wrap=True)

    add_rect(sl, Inches(0.6), H - Inches(0.58), Inches(12.1), Inches(0.03), ACCENT)
    add_text(sl,
             '"The future of solar is not just bigger panels — it is smarter panels that follow the sun, report their health, and respond to their environment automatically."',
             Inches(0.6), H - Inches(0.52), Inches(12.1), Inches(0.42),
             font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    prs.save(output_path)
    print(f"Saved: {output_path}")


# ─── DOCX ─────────────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, Inches, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

def set_cell_bg(cell, hex_color):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement('w:shd')
    shd.set(qn('w:val'),   'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'),  hex_color)
    tcPr.append(shd)

def doc_heading(doc, text, level=1, color="FFA500"):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        run.font.color.rgb = DocxRGB(int(color[0:2],16),
                                     int(color[2:4],16),
                                     int(color[4:6],16))
    return p

def doc_bullet(doc, text, style='List Bullet'):
    return doc.add_paragraph(text, style=style)

def doc_table(doc, headers, rows, col_widths=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = h
        for run in hdr_cells[i].paragraphs[0].runs:
            run.font.bold = True
        set_cell_bg(hdr_cells[i], "0D1B2A")
        for run in hdr_cells[i].paragraphs[0].runs:
            run.font.color.rgb = DocxRGB(0xFF,0xA5,0x00)
    for ri, row in enumerate(rows):
        row_cells = table.rows[ri + 1].cells
        bg = "152934" if ri % 2 == 0 else "1C3445"
        for ci, val in enumerate(row):
            row_cells[ci].text = str(val)
            set_cell_bg(row_cells[ci], bg)
            for run in row_cells[ci].paragraphs[0].runs:
                run.font.color.rgb = DocxRGB(0xCC,0xD6,0xE0)
    if col_widths:
        for i, w in enumerate(col_widths):
            for row in table.rows:
                row.cells[i].width = Inches(w)
    return table

def make_docx(output_path):
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Inches(0.9)
        section.bottom_margin = Inches(0.9)
        section.left_margin   = Inches(1.1)
        section.right_margin  = Inches(1.1)

    # ── Title ────────────────────────────────────────────────────────────────
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run("SUNFLOWER SOLAR PANEL CONTROLLER")
    run.font.size = Pt(26); run.font.bold = True
    run.font.color.rgb = DocxRGB(0xFF,0xA5,0x00)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = sub.add_run("An Intelligent, Autonomous Solar Tracking & Monitoring System")
    r2.font.size = Pt(14); r2.font.italic = True
    r2.font.color.rgb = DocxRGB(0x66,0x99,0xBB)

    sub2 = doc.add_paragraph()
    sub2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = sub2.add_run("Competition Presentation Report  |  April 2026")
    r3.font.size = Pt(11)
    r3.font.color.rgb = DocxRGB(0x88,0x99,0xAA)

    doc.add_paragraph()

    # ── 1. Introduction ──────────────────────────────────────────────────────
    doc_heading(doc, "1.  Introduction", 1)
    doc.add_paragraph(
        "The Sunflower Solar Panel Controller is a two-unit ESP32-based system that mimics the behaviour "
        "of a sunflower — deploying a solar panel at sunrise, tracking the sun throughout the day with a "
        "servo motor, and retracting the panel safely at sunset. A separate remote unit displays live sensor "
        "data on an OLED screen and allows the operator to manually override the servo and DC motor via four "
        "physical buttons, all over a private WiFi TCP connection that requires no internet infrastructure."
    )
    doc.add_paragraph(
        "The project addresses a well-known problem in solar energy: fixed panels lose 25–40% of potential "
        "energy because they cannot follow the sun. Commercial solar trackers solve this but cost ₹40,000 or "
        "more per unit. This system achieves the same core functionality for under ₹2,500 in components, "
        "making professional-grade solar tracking accessible to small farmers, off-grid homes, and developing "
        "regions worldwide."
    )

    # ── 2. System Architecture ───────────────────────────────────────────────
    doc_heading(doc, "2.  System Architecture", 1)
    doc.add_paragraph(
        "The system consists of two independent ESP32 microcontroller units that communicate over a private "
        "WiFi network using TCP sockets."
    )

    doc_heading(doc, "2.1  Solar Controller (Device 1)", 2, "2ECC71")
    doc.add_paragraph("Hardware components and GPIO assignments:")
    doc_table(doc,
        ["Component", "Interface", "GPIO Pins"],
        [
            ["L298N Motor Driver",           "Digital + PWM", "IN1=18, IN2=16, ENA=17"],
            ["Servo Motor (continuous)",     "PWM 50 Hz",     "GPIO 25"],
            ["Panel Voltage Sensor",         "ADC",           "GPIO 34"],
            ["Panel Current Sensor (ACS712)","ADC",           "GPIO 35"],
            ["Battery Voltage Sensor",       "ADC",           "GPIO 32"],
            ["RTC DS3231",                   "I²C",           "SDA=13, SCL=14"],
        ],
        [2.8, 2.0, 2.0]
    )

    doc_heading(doc, "2.2  Remote Control Unit (Device 2)", 2, "2ECC71")
    doc.add_paragraph(
        "The remote unit acts as a WiFi Access Point (192.168.4.1:5000). It hosts an OLED 128×64 dashboard "
        "showing live panel voltage, current, power, battery voltage, and the current time from the "
        "controller's RTC. Four buttons allow manual control:"
    )
    doc_table(doc,
        ["Button", "Command Sent", "Command on Release", "Effect"],
        [
            ["Button 1", "SERVO_FWD",   "SERVO_STOP", "Hold to rotate servo forward"],
            ["Button 2", "SERVO_REV",   "SERVO_STOP", "Hold to rotate servo reverse"],
            ["Button 3", "PANEL_OPEN",  "MOTOR_STOP", "Hold to open panel manually"],
            ["Button 4", "PANEL_CLOSE", "MOTOR_STOP", "Hold to close panel manually"],
        ],
        [1.2, 1.5, 1.8, 2.8]
    )

    # ── 3. Communication Protocol ────────────────────────────────────────────
    doc_heading(doc, "3.  Communication Protocol", 1)
    doc_heading(doc, "3.1  Controller → Remote (every 100 ms)", 2, "2ECC71")
    doc.add_paragraph("CSV line over TCP:")
    p = doc.add_paragraph()
    run = p.add_run("    HH:MM:SS , voltage , current , power , batteryVoltage")
    run.font.name = "Courier New"; run.font.size = Pt(11)
    doc.add_paragraph("Example:  08:42:15 , 18.34 , 2.10 , 38.51 , 12.67").runs[0].font.italic = True

    doc_heading(doc, "3.2  Remote → Controller (on button press/release)", 2, "2ECC71")
    doc_table(doc,
        ["Command", "Trigger", "Effect"],
        [
            ["SERVO_FWD",   "Hold button 1",          "Servo rotates forward — manual sun tracking"],
            ["SERVO_REV",   "Hold button 2",          "Servo rotates reverse — undo tracking angle"],
            ["SERVO_STOP",  "Release button 1 or 2",  "Servo stops immediately"],
            ["PANEL_OPEN",  "Hold button 3",          "Motor opens panel (ignored if not IDLE)"],
            ["PANEL_CLOSE", "Hold button 4",          "Motor closes panel (ignored if not IDLE)"],
            ["MOTOR_STOP",  "Release button 3 or 4",  "Motor stops immediately on release"],
            ["RESET",       "Emergency command",      "Halts everything, clears all state flags"],
        ],
        [1.8, 2.2, 3.2]
    )

    # ── 4. Key Technical Features ────────────────────────────────────────────
    doc_heading(doc, "4.  Key Technical Features", 1)

    doc_heading(doc, "4.1  Non-Blocking Architecture", 2, "2ECC71")
    for b in [
        "Main loop runs every 10 ms — never uses delay(). All operations use millis() timers.",
        "Motor, servo, WiFi reconnect, and sensors all run concurrently without pausing each other.",
        "WiFi client.connect() is only attempted when motorState == MOTOR_IDLE to prevent blocking the PWM ramp (connect() can stall the CPU for up to 500ms).",
        "RTC is read over I²C every 1 second and cached — avoids 100× I²C overhead on every 10ms loop iteration.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "4.2  Motor State Machine", 2, "2ECC71")
    for b in [
        "States: MOTOR_IDLE → MOTOR_OPENING / MOTOR_CLOSING / MOTOR_MANUAL → MOTOR_IDLE",
        "New commands (schedule or manual) are only accepted when motorState == MOTOR_IDLE, preventing double-trigger.",
        "The opened/closed flags are set after the motor physically finishes its run — not when the command fires — reflecting real-world panel position.",
        "Emergency RESET instantly calls hardBrake(), sets motorState = MOTOR_IDLE, and clears all timing variables.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "4.3  PWM Soft-Ramp Motor Control", 2, "2ECC71")
    for b in [
        "Custom MotorL298N library ramps PWM from 0% → 100% over approximately 5 seconds (1% per 50ms by default).",
        "A 30% kickstart pulse is applied immediately on setSpeed() to overcome static friction before ramping.",
        "hardBrake() instantly cuts all PWM and direction pins — used for emergency RESET.",
        "Smooth ramp significantly reduces mechanical shock, extending the life of the gearbox, belt, and panel mounting hardware.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "4.4  Automatic Current Sensor Calibration", 2, "2ECC71")
    for b in [
        "On every boot, 500 ADC samples are averaged over 1 second with no load connected.",
        "The resulting mean voltage is stored as ACS_OFFSET — the zero-current reference.",
        "This eliminates offset drift caused by temperature changes and supply voltage variation.",
        "No manual potentiometer or trim adjustment is required — the system self-configures each startup.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "4.5  Servo Sun Tracking", 2, "2ECC71")
    for b in [
        "Auto-tracking fires every 15 minutes while the panel is open (configurable via SERVO_TRACK_INTERVAL).",
        "Duration is mathematically calculated from SERVO_RPM and SERVO_TRACK_REVS — not guessed.",
        "Priority order: Manual command > Auto-track > Stop. Manual override pauses auto-tracking instantly.",
        "Auto-tracking is cancelled immediately when the panel closes at end of day.",
    ]: doc_bullet(doc, b)

    # ── 5. Large Scale Implementation ────────────────────────────────────────
    doc_heading(doc, "5.  Large Scale Implementation", 1)

    doc_heading(doc, "5.1  Hardware Scaling", 2, "2ECC71")
    for b in [
        "Replace the single ESP32 controller with a Raspberry Pi or industrial PLC as a central hub managing hundreds of panel units over a network.",
        "The L298N motor driver is replaced with 3-phase motor controllers for the larger DC or AC actuators used in commercial trackers.",
        "Stepper motors with quadrature encoders replace the continuous-rotation servo, providing precise angular position control and eliminating open-loop drift.",
        "Industrial-grade CT clamps rated for 100A+ replace the ACS712 modules for high-current installations.",
        "The OLED display is replaced by a web-based Grafana or React dashboard showing fleet-wide metrics on a map, with per-panel status indicators.",
        "Dual-axis tracking: a second motor added for elevation angle control increases annual yield by a further 10–15% beyond single-axis.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "5.2  Network and Software Scaling", 2, "2ECC71")
    for b in [
        "Replace direct TCP with an MQTT broker (Mosquitto). Each panel publishes sensor data to its own topic; a central server subscribes to all topics simultaneously.",
        "OTA (Over-the-Air) firmware updates push new code to all panels in a field without physical access to each unit.",
        "Time-series database (InfluxDB or TimescaleDB) stores all sensor readings for historical analysis, fault detection, and energy yield reporting.",
        "A REST API allows any monitoring tool (mobile app, web browser, SCADA system) to query panel status, issue commands, and update schedules.",
        "GPS-based sun angle calculation using the NOAA solar position algorithm computes optimal panel angle from latitude, longitude, and UTC time — eliminating the fixed 15-minute interval assumption.",
        "Cloud backup with automated daily energy yield reports delivered via email or SMS to the owner.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "5.3  Physical Deployment", 2, "2ECC71")
    for b in [
        "Panels mounted on single-axis tracker rails (east-west rotation) driven by the DC motor — the most common commercial solar tracker design.",
        "Wind speed sensor auto-retracts all panels when gusts exceed a configurable safe threshold, preventing storm damage.",
        "IP65-rated weatherproof enclosures for all electronics with cable glands; surge protection and reverse-polarity protection on all power inputs.",
        "Modular wiring harness — each panel unit is identical and plug-and-play, reducing installation time and training requirements.",
        "Staggered motor start schedules prevent current surges when all 100+ panels attempt to open simultaneously at sunrise.",
    ]: doc_bullet(doc, b)

    # ── 6. System Improvements ───────────────────────────────────────────────
    doc_heading(doc, "6.  System Improvements", 1)

    doc_heading(doc, "6.1  Precision Sun Tracking", 2, "2ECC71")
    for b in [
        "Add an LDR (light-dependent resistor) differential pair — the servo rotates until both sensors read equal illumination, meaning the panel is facing the sun directly. Replaces the open-loop timed pulse.",
        "Closed-loop position control: rotary encoder on the servo shaft confirms actual angular position and enables the controller to command precise angles rather than timed rotations.",
        "MPPT algorithm (Maximum Power Point Tracking): continuously adjust panel angle to maximise real-time power output, not just geometric alignment with the sun. Accounts for cloud diffusion and partial shading.",
        "Add an elevation axis motor for full 2-axis tracking — adds 10–15% annual energy yield over single-axis alone.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "6.2  Reliability and Safety", 2, "2ECC71")
    for b in [
        "Motor stall detection: add a dedicated current sensor on the motor supply line. A current spike without movement triggers an immediate halt and fault alert.",
        "Panel position encoder: know the exact open angle at all times. Prevents over-driving the motor if the panel is mechanically obstructed.",
        "ESP32 hardware watchdog timer: auto-reboots the system if the firmware hangs, ensuring the panel is never left in an unknown state.",
        "Redundant time source: sync from NTP when WiFi is available; fall back to DS3231 RTC when offline. RTC drift corrected automatically after each NTP sync.",
        "Surge protection and TVS diodes on all ADC input lines to prevent sensor damage from lightning or motor back-EMF.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "6.3  Energy Intelligence", 2, "2ECC71")
    for b in [
        "State-of-charge calculation: derive battery percentage from the discharge voltage curve rather than relying on raw voltage alone.",
        "Watt-hour accumulation: log daily energy yield (kWh) to flash memory for comparison against historical and predicted values.",
        "Low battery protection: automatically retract the panel and enter sleep mode when battery voltage drops critically low, protecting the battery from deep discharge.",
        "Efficiency monitoring: compare actual yield to expected yield based on solar irradiance models; alert when efficiency drops more than 10% below threshold (indicating soiling or shading).",
        "Panel temperature sensor: derate power estimate on hot days (approximately 0.3% per °C above 25°C for silicon cells).",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "6.4  User Interface", 2, "2ECC71")
    for b in [
        "Replace the OLED with a colour TFT touchscreen showing scrollable voltage, current, and power graphs with time-axis.",
        "Mobile app via Bluetooth or WiFi for configuration: set open/close times, tracking interval, and RPM calibration without touching firmware.",
        "Push notifications when power drops below expected threshold or a fault is detected — delivered via SMS, WhatsApp, or email.",
        "Web interface with historical charts, exportable CSV data, and remote schedule configuration accessible from any browser.",
        "Voice alerts via a small buzzer: distinct tones for open complete, close complete, and fault events.",
    ]: doc_bullet(doc, b)

    # ── 7. Future Benefits ───────────────────────────────────────────────────
    doc_heading(doc, "7.  Future Benefits and Applications", 1)

    doc_heading(doc, "7.1  Energy Independence", 2, "2ECC71")
    for b in [
        "Single-axis tracking increases annual energy yield by 25–40% over a fixed panel installation.",
        "Adding a second (elevation) axis increases yield by a further 10–15%, for a total gain of up to 55% over fixed.",
        "More energy from the same number of panels reduces the total hardware cost and land area required for a given power target.",
        "Longer battery cycle life: the tracker maximises energy extraction during peak sun hours, reducing the number of days the battery must supply energy without solar recharge.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "7.2  Agrivoltaics (Agriculture + Solar)", 2, "2ECC71")
    for b in [
        "Tilted solar panels allow morning and evening light to reach crops underneath while blocking harsh midday sun.",
        "Crops grown under tracked panels benefit from reduced heat stress, with studies showing up to 60% less water consumption in some climates.",
        "Dual income stream: farmers generate and sell electricity while continuing to grow crops on the same land.",
        "This system's automated tilt angle control makes it a ready-to-use foundation for agrivoltaic installations at the village level.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "7.3  Smart Grid Integration", 2, "2ECC71")
    for b in [
        "Real-time power data at 100ms resolution feeds directly into grid balancing and demand-response systems.",
        "When panels report surplus generation, smart switches can automatically divert energy to water heating, EV charging, or grid export.",
        "Predictive generation forecasting: historical tracking data combined with weather APIs helps grid operators plan backup generation capacity.",
        "Micro-grid: a cluster of these units with a shared battery bank can power an entire village with automatic load balancing.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "7.4  Developing World and Disaster Relief", 2, "2ECC71")
    for b in [
        "No internet required — the two ESP32s form their own private WiFi network, deployable in remote areas with zero infrastructure.",
        "Total BOM cost under ₹2,500 (~$30) makes professional solar tracking accessible to smallholder farmers globally.",
        "Can be assembled with basic tools and locally available modules — no specialist manufacturing required.",
        "Portable solar tracker + battery provides uninterrupted power for medical equipment, water pumps, and communications in disaster zones.",
    ]: doc_bullet(doc, b)

    doc_heading(doc, "7.5  Environmental Impact", 2, "2ECC71")
    impact_rows = [
        ["1 rooftop panel (300W)",  "+~100 Wh/day",   "~15 kg CO₂/year"],
        ["100-panel farm (30 kW)",  "+~10 kWh/day",   "~1,500 kg CO₂/year"],
        ["Utility-scale (1 MW)",    "+~350 kWh/day",  "~50,000 kg CO₂/year"],
    ]
    doc_table(doc,
        ["Scale", "Daily Energy Gain vs Fixed", "Annual CO₂ Offset"],
        impact_rows, [3.0, 2.5, 2.5]
    )
    doc.add_paragraph(
        "Based on 35% improvement from single-axis tracking, 300 sunny days/year, and 0.4 kg CO₂/kWh grid average (IEA 2024)."
    ).runs[0].font.italic = True if doc.paragraphs[-1].runs else None

    # ── 8. Competitive Advantages ────────────────────────────────────────────
    doc_heading(doc, "8.  Competitive Advantages", 1)
    advantages = [
        ("No Internet Required",
         "The private WiFi network between the two ESP32s works in fields, deserts, and rooftops without any infrastructure. Unlike cloud-dependent systems, it cannot be disrupted by a server outage."),
        ("Custom PWM Motor Library",
         "Soft-ramp 0→90% over 5 seconds with kickstart prevents gear shock, motor burnout, and belt snap. No commercial tracker at this price offers configurable ramp control."),
        ("Self-Calibrating Sensors",
         "ACS712 zero offset auto-calibrated on every boot with 500 samples — compensates for temperature drift and eliminates manual potentiometer tuning."),
        ("State-Machine Safety",
         "Motor only accepts new commands when IDLE. Flags reflect physical panel state, not command state. Single RESET halts everything to a known safe state."),
        ("100ms Live Data Stream",
         "10 readings per second is sufficient to detect transient shadows, partial shading, and wiring faults — far beyond the 1-second resolution of most hobby systems."),
        ("Extensible Protocol",
         "Adding a new remote command requires one line in the buttons[] array. No other code changes needed — engineers can extend the system in minutes."),
        ("Under ₹2,500 Total Cost",
         "Commercial solar trackers start at ₹40,000+. This system delivers the core functionality for 94% less, making it viable for mass deployment in developing markets."),
    ]
    for title, body in advantages:
        p = doc.add_paragraph()
        run_t = p.add_run(f"{title}: ")
        run_t.font.bold = True
        run_t.font.color.rgb = DocxRGB(0xFF,0xA5,0x00)
        run_b = p.add_run(body)
        run_b.font.color.rgb = DocxRGB(0x33,0x33,0x33)

    # ── 9. Bill of Materials ─────────────────────────────────────────────────
    doc_heading(doc, "9.  Bill of Materials", 1)
    doc_table(doc,
        ["Component", "Qty", "Unit Price", "Total", "Purpose"],
        [
            ["ESP32 Dev Board",           "2", "₹400", "₹800",  "Controller + Remote unit"],
            ["L298N Motor Driver",        "1", "₹150", "₹150",  "DC motor control"],
            ["DC Gear Motor (12V)",       "1", "₹350", "₹350",  "Panel open/close"],
            ["Continuous Rotation Servo", "1", "₹280", "₹280",  "Sun tracking"],
            ["ACS712-30A Current Sensor", "1", "₹120", "₹120",  "Panel current"],
            ["Voltage Divider Module",    "2", "₹45",  "₹90",   "Panel + battery voltage"],
            ["DS3231 RTC Module",         "1", "₹180", "₹180",  "Schedule timekeeping"],
            ["SSD1306 OLED 128×64",       "1", "₹220", "₹220",  "Live dashboard display"],
            ["Push Buttons (pack)",       "1", "₹50",  "₹50",   "Manual control buttons"],
            ["Misc (wires, PCB)",         "1", "₹200", "₹200",  "Wiring and prototyping"],
            ["TOTAL",                     "",  "",     "₹2,440",""],
        ],
        [2.8, 0.5, 1.0, 0.8, 2.5]
    )

    # ── 10. Conclusion ───────────────────────────────────────────────────────
    doc_heading(doc, "10.  Conclusion", 1)
    doc.add_paragraph(
        "The Sunflower Solar Panel Controller demonstrates that intelligent, autonomous solar tracking is not "
        "a luxury reserved for large utility companies. With under ₹2,500 of components and carefully engineered "
        "firmware, it is possible to build a system that:"
    )
    for b in [
        "Increases annual solar energy yield by 25–40% through automated sun tracking.",
        "Monitors real-time panel voltage, current, and power at 100ms resolution.",
        "Provides a wireless OLED dashboard with live data and manual override — no internet required.",
        "Uses a robust state-machine architecture that is safe, non-blocking, and extensible.",
        "Scales directly to large solar farms via MQTT, cloud databases, and OTA firmware updates.",
        "Saves up to 50,000 kg of CO₂ per year per MW of solar capacity compared to fixed installations.",
    ]: doc_bullet(doc, b)

    doc.add_paragraph()
    quote = doc.add_paragraph()
    quote.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = quote.add_run(
        '"The future of solar is not just bigger panels — it is smarter panels that follow the sun, '
        'report their health, and respond to their environment automatically."'
    )
    r.font.italic = True
    r.font.color.rgb = DocxRGB(0x55,0x77,0x99)

    doc.save(output_path)
    print(f"Saved: {output_path}")


# ─── Main ─────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    base = "/home/pavan/Downloads/SOLAR_CONREOLLER"
    make_pptx(f"{base}/Sunflower_Solar_Controller_Presentation.pptx")
    make_docx(f"{base}/Sunflower_Solar_Controller_Report.docx")
    print("Done.")
